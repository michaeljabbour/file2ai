import os
from pptx import Presentation


def create_test_presentation():
    """Create a test PowerPoint presentation with sample content."""
    prs = Presentation()

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created by file2ai"

    # Add a content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Sample Content"
    tf = body.text_frame
    tf.text = "This is a test slide"
    p = tf.add_paragraph()
    p.text = "• Bullet point 1"
    p = tf.add_paragraph()
    p.text = "• Bullet point 2"

    # Save the presentation
    # Ensure we're using absolute paths from project root
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    exports_dir = os.path.join(project_root, "exports")

    # Create exports directory if it doesn't exist
    os.makedirs(exports_dir, exist_ok=True)

    output_path = os.path.join(exports_dir, "test.pptx")
    prs.save(output_path)
    print(f"Created test presentation at: {output_path}")


if __name__ == "__main__":
    create_test_presentation()
