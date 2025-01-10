from pathlib import Path
import docx
from docx.shared import Inches
import openpyxl
from pptx import Presentation
from pptx.util import Inches as PptxInches

def create_test_files():
    test_dir = Path("test_files")
    test_dir.mkdir(exist_ok=True)
    
    # Create HTML
    html_path = test_dir / "test.html"
    html_content = """<!DOCTYPE html>
<html><body><h1>Test Document</h1>
<p>This is a test document created for file2ai testing.</p>
</body></html>"""
    html_path.write_text(html_content)
    
    # Create DOCX
    doc = docx.Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test document created for file2ai testing.')
    doc.save(str(test_dir / "test.docx"))
    
    # Create XLSX
    wb = openpyxl.Workbook()
    ws = wb.create_sheet("Sheet1")
    ws.title = "Sheet1"
    ws.cell(row=1, column=1, value='Test')
    ws.cell(row=1, column=2, value='Document')
    ws.cell(row=2, column=1, value='This is')
    ws.cell(row=2, column=2, value='a test.')
    wb.save(str(test_dir / "test.xlsx"))
    
    # Create PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created for file2ai testing"
    prs.save(str(test_dir / "test.pptx"))
    
    # Create PDF file using reportlab
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    pdf_path = test_dir / "test.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    c.drawString(100, 750, "Test Document")
    c.drawString(100, 700, "This is a test file created for file2ai testing.")
    c.save()

if __name__ == "__main__":
    create_test_files()
