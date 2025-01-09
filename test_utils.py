"""Utilities for test script progress tracking and file creation."""
import sys
import time
import logging
from tqdm import tqdm
import subprocess
from typing import List, Optional
from threading import Lock
from pathlib import Path
import docx
from docx.shared import Inches

# Mock mimetypes before importing openpyxl
import unittest.mock
import sys
import types
from pathlib import Path

# Create a mutable module class
class MutableModule(types.ModuleType):
    def __init__(self, name):
        super().__init__(name)
        self._namespace = {}
    
    def __setattr__(self, name, value):
        if name == '_namespace':
            super().__setattr__(name, value)
        else:
            self._namespace[name] = value
    
    def __getattr__(self, name):
        try:
            return self._namespace[name]
        except KeyError:
            raise AttributeError(f"module '{self.__name__}' has no attribute '{name}'")

# Create base mime types implementation
class MimeTypes:
    def __init__(self):
        self.types_map = {
            '.html': 'text/html',
            '.htm': 'text/html',
            '.txt': 'text/plain',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.xls': 'application/vnd.ms-excel',
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.xml': 'application/xml'
        }
        self.common_types = self.types_map.copy()
        
    def guess_type(self, url, strict=True):
        ext = Path(url).suffix.lower()
        return self.types_map.get(ext, 'application/octet-stream'), None
        
    def add_type(self, type, ext, strict=True):
        if ext[0] != '.':
            ext = '.' + ext
        ext = ext.lower()
        self.types_map[ext] = type
        self.common_types[ext] = type
        
    def read(self, filename):
        pass
        
    def readfp(self, fp, strict=True):
        pass

# Create mock module
mock_mimetypes = MutableModule('mimetypes')
mime_types_instance = MimeTypes()

# Set up module attributes and methods
mock_mimetypes.MimeTypes = MimeTypes
mock_mimetypes.guess_type = mime_types_instance.guess_type
mock_mimetypes.add_type = mime_types_instance.add_type
mock_mimetypes.init = unittest.mock.MagicMock()
mock_mimetypes.knownfiles = []
mock_mimetypes.inited = True
mock_mimetypes.suffix_map = {}
mock_mimetypes.encodings_map = {}
mock_mimetypes.MIME_TYPES_ENCODINGS = 1
mock_mimetypes.MIME_TYPES_TYPES = 2

# Replace real mimetypes module with our mock
sys.modules['mimetypes'] = mock_mimetypes

import openpyxl
from pptx import Presentation
from pptx.util import Inches as PptxInches
from weasyprint import HTML

logger = logging.getLogger(__name__)

# Global progress bar for overall progress
_overall_progress = None
_progress_lock = Lock()

def init_overall_progress(total_steps: int = 10):
    """Initialize the overall progress bar that stays at the top."""
    global _overall_progress
    with _progress_lock:
        if _overall_progress is None:
            _overall_progress = tqdm(total=total_steps, 
                                  desc="Overall Progress",
                                  position=0,
                                  leave=True,
                                  bar_format='{desc}: {percentage:3.0f}%|{bar}| {n_fmt}/{total_fmt}')

def run_with_progress(command: List[str], desc: str, total: Optional[int] = None) -> int:
    """Run a command with a progress bar."""
    process = subprocess.Popen(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        universal_newlines=True,
    )
    
    # Use simple progress indicator
    with tqdm(desc=desc, 
             total=1,
             position=1,
             leave=False) as pbar:
        while process.poll() is None:
            time.sleep(0.1)
        pbar.update(1)
    
    # Update overall progress
    global _overall_progress
    with _progress_lock:
        if _overall_progress:
            _overall_progress.update(1)
    
    return process.returncode

def show_spinner(desc: str, duration: float) -> None:
    """Show a spinner for a fixed duration."""
    with tqdm(desc=desc,
             total=1,
             position=1,
             leave=False) as pbar:
        time.sleep(duration)
        pbar.update(1)

def create_test_files(test_dir: str) -> None:
    """Create test files with progress tracking."""
    test_path = Path(test_dir)
    test_path.mkdir(parents=True, exist_ok=True)
    
    with tqdm(total=5,
             desc="Creating test files",
             position=1,
             leave=False,
             bar_format='{desc}: {percentage:3.0f}%|{bar}| {n_fmt}/{total_fmt}') as pbar:
        
        # Create HTML and PDF
        html_path = test_path / "test.html"
        html_content = """<!DOCTYPE html>
<html><body><h1>Test Document</h1>
<p>This is a test document created for file2ai testing.</p>
</body></html>"""
        html_path.write_text(html_content)
        HTML(string=html_content).write_pdf(str(test_path / "test.pdf"))
        pbar.update(1)
        
        # Create DOCX
        doc = docx.Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test document created for file2ai testing.')
        doc.save(str(test_path / "test.docx"))
        pbar.update(1)
        
        # Create XLSX
        wb = openpyxl.Workbook()
        if wb.active is not None:
            ws = wb.active
            ws['A1'] = 'Test'
            ws['B1'] = 'Document'
            ws['A2'] = 'This is'
            ws['B2'] = 'a test.'
            wb.save(str(test_path / "test.xlsx"))
        pbar.update(1)
        
        # Create PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title is not None:
            slide.shapes.title.text = "Test Presentation"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Created for file2ai testing"
        prs.save(str(test_path / "test.pptx"))
        pbar.update(1)
        
        # Create simple text file
        text_path = test_path / "test.txt"
        text_path.write_text("This is a test file created for file2ai testing.\n")
        pbar.update(1)
    
    # Update overall progress
    global _overall_progress
    with _progress_lock: 
        if _overall_progress:
            _overall_progress.update(1)

def install_deps_with_progress(packages: List[str]) -> None:
    """Install pip packages with a progress bar."""
    with tqdm(total=len(packages),
             desc="Installing dependencies",
             position=1,
             leave=False,
             bar_format='{desc}: {percentage:3.0f}%|{bar}| {n_fmt}/{total_fmt}') as pbar:
        for package in packages:
            if package == "--upgrade pip":
                subprocess.run([sys.executable, "-m", "pip", "install", "--upgrade", "pip"], check=True)
            elif package == "-e .[test,web]":
                subprocess.run([sys.executable, "-m", "pip", "install", "-e", ".[test,web]"], check=True)
            else:
                logger.warning(f"Unexpected package installation request: {package}")
            pbar.update(1)
    
    # Update overall progress
    global _overall_progress
    with _progress_lock:
        if _overall_progress:
            _overall_progress.update(1)
